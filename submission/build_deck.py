#!/usr/bin/env python3
"""Generate the Gauntlet pitch deck (6 slides) as a styled PowerPoint.
Dark security-console aesthetic to match the product. Run: python3 submission/build_deck.py
Output: submission/Gauntlet-Deck.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK = RGBColor(0x0E, 0x11, 0x16)
SURFACE = RGBColor(0x16, 0x1B, 0x22)
EDGE = RGBColor(0x2B, 0x33, 0x3D)
TEXT = RGBColor(0xF2, 0xF4, 0xF7)
MUTED = RGBColor(0x9A, 0xA4, 0xAF)
ACCENT = RGBColor(0x4C, 0xC2, 0xFF)
SIGNAL = RGBColor(0x3D, 0xDC, 0x97)
ALERT = RGBColor(0xF2, 0x55, 0x5A)
WARN = RGBColor(0xF2, 0xC1, 0x4E)

SANS = "Arial"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height


def bg(slide):
    r = slide.shapes.add_shape(1, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = INK
    r.line.fill.background()
    r.shadow.inherit = False
    return r


def bar(slide):
    b = slide.shapes.add_shape(1, 0, 0, Inches(0.16), H)
    b.fill.solid(); b.fill.fore_color.rgb = ACCENT
    b.line.fill.background(); b.shadow.inherit = False


def tag(slide, n):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run(); run.text = "GAUNTLET · autonomous AI red-team"
    run.font.name = MONO; run.font.size = Pt(10); run.font.color.rgb = MUTED
    r2 = p.add_run(); r2.text = f"      {n}/6"
    r2.font.name = MONO; r2.font.size = Pt(10); r2.font.color.rgb = EDGE


def textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    return tf


def add_line(tf, text, size, color, font=SANS, bold=False, space_before=6, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(space_before); p.space_after = Pt(2)
    run = p.add_run(); run.text = text
    run.font.name = font; run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color
    return p


def kicker(slide, text):
    tf = textbox(slide, Inches(0.6), Inches(0.55), Inches(12), Inches(0.5))
    add_line(tf, text, 14, ACCENT, font=MONO, bold=True, first=True)


def title(slide, text, color=TEXT):
    tf = textbox(slide, Inches(0.6), Inches(1.0), Inches(12.1), Inches(1.4))
    add_line(tf, text, 36, color, bold=True, first=True)


def bullets(slide, items, top=2.6):
    tf = textbox(slide, Inches(0.7), Inches(top), Inches(12), Inches(4.0))
    for i, (txt, col) in enumerate(items):
        p = add_line(tf, "•  " + txt, 18, col, space_before=10, first=(i == 0))


# Slide 1 — title / problem
s = prs.slides.add_slide(BLANK); bg(s); bar(s)
kicker(s, "BEYOND TOMORROW HACKATHON")
tf = textbox(s, Inches(0.6), Inches(2.2), Inches(12), Inches(2.2))
add_line(tf, "Everyone is shipping AI.", 44, TEXT, bold=True, first=True)
add_line(tf, "Almost no one is testing it for attacks.", 44, MUTED, bold=True, space_before=2)
tf2 = textbox(s, Inches(0.7), Inches(4.7), Inches(12), Inches(2))
add_line(tf2, "Prompt injection is the #1 risk in the OWASP LLM Top 10 (2025), and it takes nothing but text.", 18, TEXT, first=True)
add_line(tf2, "A leaked system prompt or an exfiltrated record is a real, current, expensive failure.", 18, MUTED)
add_line(tf2, "Throw your AI in. See what survives.", 20, ACCENT, font=MONO, bold=True, space_before=16)
tag(s, 1)

# Slide 2 — what it is
s = prs.slides.add_slide(BLANK); bg(s); bar(s)
kicker(s, "WHAT GAUNTLET IS")
title(s, "An autonomous red-team that attacks, scores, and fixes your AI in one loop.")
bullets(s, [
    ("Point it at a bot, a pasted system prompt, a real model, or your own HTTP endpoint.", TEXT),
    ("A generative, target-aware attacker fires OWASP-mapped probes, live.", TEXT),
    ("One click installs a runtime guard and re-tests. The grade climbs F to A.", TEXT),
    ("Built for the person shipping the feature, not just security engineers.", MUTED),
])
tf = textbox(s, Inches(0.7), Inches(5.7), Inches(12), Inches(0.8))
add_line(tf, "attack  ->  score  ->  one-click guard  ->  rescore", 22, ACCENT, font=MONO, bold=True, first=True)
tag(s, 2)

# Slide 3 — before/after
s = prs.slides.add_slide(BLANK); bg(s); bar(s)
kicker(s, "THE WOW: BEFORE AND AFTER")
title(s, "F to A, on the real app, in two clicks.")
# Two cards
card1 = s.shapes.add_shape(1, Inches(0.9), Inches(2.7), Inches(5.4), Inches(3.2))
card1.fill.solid(); card1.fill.fore_color.rgb = SURFACE; card1.line.color.rgb = ALERT; card1.line.width = Pt(1.5); card1.shadow.inherit = False
card2 = s.shapes.add_shape(1, Inches(7.0), Inches(2.7), Inches(5.4), Inches(3.2))
card2.fill.solid(); card2.fill.fore_color.rgb = SURFACE; card2.line.color.rgb = SIGNAL; card2.line.width = Pt(1.5); card2.shadow.inherit = False
for shape, grade, gcol, sub in [(card1, "F", ALERT, "SupportBot: 10 of 13 attacks got in"), (card2, "A", SIGNAL, "After the guard: 0 got through")]:
    tf = shape.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = grade; r.font.name = MONO; r.font.size = Pt(96); r.font.bold = True; r.font.color.rgb = gcol
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = sub; r2.font.name = SANS; r2.font.size = Pt(14); r2.font.color.rgb = MUTED
tf = textbox(s, Inches(0.7), Inches(6.1), Inches(12), Inches(0.7))
add_line(tf, "The remediation panel names what changed for each attack. Reproducible: npm run eval.", 16, MUTED, first=True)
tag(s, 3)

# Slide 4 — how it works
s = prs.slides.add_slide(BLANK); bg(s); bar(s)
kicker(s, "HOW IT WORKS")
title(s, "A streaming, multi-stage agent loop.")
bullets(s, [
    ("Planner picks OWASP families. Attacker generates probes (incl. multi-turn and indirect).", TEXT),
    ("Target adapter responds. A layered guard filters input and sanitizes output.", TEXT),
    ("A canary oracle decides compromise, refusal-gated to avoid false positives.", TEXT),
    ("A scorer maps findings to the OWASP LLM Top 10. Everything streams over SSE.", TEXT),
])
tf = textbox(s, Inches(0.7), Inches(5.7), Inches(12), Inches(0.9))
add_line(tf, "Next.js 16 · React 19 · TypeScript · Tailwind · Vercel · Vitest · Playwright", 16, ACCENT, font=MONO, bold=True, first=True)
tag(s, 4)

# Slide 5 — measured (honest)
s = prs.slides.add_slide(BLANK); bg(s); bar(s)
kicker(s, "WE MEASURED IT (THE HONEST SLIDE)")
title(s, "Numbers we earned, not adjectives.")
bullets(s, [
    ("Oracle: 0% false positives, 11% false negatives on an 18-case labeled set (the miss is documented).", TEXT),
    ("Reproducible grades: SupportBot F(10/13), DevAssistant F(5/13), PolicyBot F(9/13), all to A.", TEXT),
    ("A real frontier model mostly holds: ~40 probes, only conflicting-instruction prompts leaked (~1 in 6).", TEXT),
    ("The real risk is prompt misconfiguration, and that is what Gauntlet surfaces for your app.", SIGNAL),
])
tag(s, 5)

# Slide 6 — impact + ask
s = prs.slides.add_slide(BLANK); bg(s); bar(s)
kicker(s, "IMPACT AND THE ASK")
title(s, "Adversarial testing any builder can run.")
bullets(s, [
    ("The guard is free and deterministic at runtime; the production path is a local classifier.", TEXT),
    ("Scales to any app via bring-your-own endpoint and a one-line CI check (next).", TEXT),
    ("Live at gauntlet-seven.vercel.app, open source on GitHub.", TEXT),
])
tf = textbox(s, Inches(0.7), Inches(5.4), Inches(12), Inches(1.2))
add_line(tf, "Throw your AI in. See what survives.", 28, ACCENT, font=MONO, bold=True, first=True)
tag(s, 6)

import os
out = os.path.join(os.path.dirname(__file__), "Gauntlet-Deck.pptx")
prs.save(out)
print("wrote", out)
